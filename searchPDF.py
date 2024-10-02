import os
from pptx import Presentation

def search_pptx_for_string(directory, search_string):
    results = []
    
    # Walk through all subdirectories and files
    for foldername, subfolders, filenames in os.walk(directory):
        for filename in filenames:
            if filename.endswith('.pptx'):
                file_path = os.path.join(foldername, filename)
                try:
                    # Load the PowerPoint file
                    presentation = Presentation(file_path)
                    
                    # Search each slide for the string
                    for slide in presentation.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                if search_string in shape.text:
                                    results.append(file_path)
                                    break  # Found in this file, no need to search more slides in this file
                except Exception as e:
                    print(f"Error reading {file_path}: {e}")
    
    return results

# Example usage
directory_to_search = '.'  # Replace with your directory path
search_string = 'SIGINT'

matching_files = search_pptx_for_string(directory_to_search, search_string)
if matching_files:
    print("Found in the following files:")
    for file in matching_files:
        print(file)
else:
    print("No matches found.")